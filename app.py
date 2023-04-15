from flask import Flask, render_template, request, send_file
from pptx import Presentation
from flask import Flask, render_template, url_for

from pptx.util import Inches

app = Flask(__name__)


@app.route('/')
def index():
    image_url = url_for('static', filename='images/chiki.gif')
    return render_template('land.html', image_url=image_url)


@app.route('/form')
def form():
    return render_template('input_form.html')


@app.route('/submit', methods=['POST'])
def generate_ppt():
    # Get user input from the HTML form
    title_text = request.form['project_title']
    members_text = request.form['project_members']
    about_text = request.form['project_description']
    abstract_text = request.form['project-abstract']
    Solution_text = request.form['project-solution']
    Techstack_text = request.form['project-techstack']

    # image_path = request.form['project_image']
    # Create a new PowerPoint presentation
    pr = Presentation()

    # Slide 1 - Title slide
    slide1_layout = pr.slide_layouts[1]
    slide1 = pr.slides.add_slide(slide1_layout)
    title_shape = slide1.shapes.title
    title_shape.text = "Introduction"
    bullet_shape = slide1.placeholders[1]
    bullet_shape.text = members_text

    # Slide 2 - About slide
    slide2_layout = pr.slide_layouts[1]
    slide2 = pr.slides.add_slide(slide2_layout)
    title_shape = slide2.shapes.title
    title_shape.text = 'Problem Statement'
    bullet_shape = slide2.placeholders[1]
    bullet_shape.text = about_text

    # slide 3
    slide3_layout = pr.slide_layouts[1]
    slide3 = pr.slides.add_slide(slide3_layout)
    title_shape = slide3.shapes.title
    title_shape.text = 'Abstract'
    bullet_shape = slide3.placeholders[1]
    bullet_shape.text = abstract_text

    # slide4
    slide4_layout = pr.slide_layouts[1]
    slide4 = pr.slides.add_slide(slide4_layout)
    title_shape = slide4.shapes.title
    title_shape.text = 'Solution'
    bullet_shape = slide4.placeholders[1]
    bullet_shape.text = Solution_text

    # slide5
    slide5_layout = pr.slide_layouts[1]
    slide5 = pr.slides.add_slide(slide5_layout)
    title_shape = slide5.shapes.title
    title_shape.text = 'Tech-Stack Used'
    bullet_shape = slide5.placeholders[1]
    bullet_shape.text = Techstack_text

    # # Slide 3 - Image slide
    # slide3_layout = pr.slide_layouts[5]
    # slide3 = pr.slides.add_slide(slide3_layout)
    # title_shape = slide3.shapes.title
    # title_shape.text = 'Image'
    # image_shape = slide3.shapes.add_picture(image_path, Inches(1), Inches(1), height=Inches(5), width=Inches(5))

    # Save the presentation
    file_name = request.form['file_name']
    pr.save(file_name + '.pptx')

    return render_template('success.html', file_name=file_name)

@app.route('/download/<string:file_name>', methods=['GET'])
def download_file(file_name):
    try:
        return send_file(file_name+'.pptx', as_attachment=True)
    except Exception as e:
        return str(e)



if __name__ == '__main__':
    app.run(debug=True)
